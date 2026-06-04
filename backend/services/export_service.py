"""
TriDoc 导出服务 — 将翻译/润色结果回写到原始文件格式。

支持的导出:
  - 翻译版: 单一语言替换
  - 对照版: 原文 + 译文逐段对照（PDF 双栏，PPTX/Word 交替段落）
"""

from __future__ import annotations

import logging
from io import BytesIO
from pathlib import Path

logger = logging.getLogger("tri_doc.export")

BASE_DIR = Path(__file__).resolve().parent.parent
UPLOAD_DIR = BASE_DIR / "uploads"

# ============================================================================
# 统一导出入口
# ============================================================================
def export_document(file_id: str, target_lang: str, mode: str = "translated") -> BytesIO:
    """
    导出文件。

    Args:
        file_id: 原始文件 ID
        target_lang: 目标语言 (ja/en/zh)
        mode: "translated" | "bilingual" (对照版)

    Returns:
        BytesIO: 导出文件流
    """
    ext = Path(file_id).suffix.lower()
    if ext == ".pdf":
        return _export_pdf(file_id, target_lang, mode)
    elif ext == ".pptx":
        return _export_pptx(file_id, target_lang, mode)
    elif ext == ".docx":
        return _export_word(file_id, target_lang, mode)
    else:
        raise ValueError(f"Unsupported export format: {ext}")


def _load_translated_pages(file_id: str, target_lang: str) -> list[dict]:
    """加载翻译/润色后的页面数据（从 sidecar JSON）。"""
    import json
    sidecar_path = UPLOAD_DIR / f"{file_id}.json"
    if not sidecar_path.exists():
        raise FileNotFoundError(f"Sidecar not found: {sidecar_path}")

    with open(sidecar_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    # 优先级: final > aligned > translated
    for key in ["final", "aligned", "translations"]:
        if key in data and target_lang in data[key]:
            return data[key][target_lang]

    raise ValueError(f"No translation found for language: {target_lang}")


# ============================================================================
# PDF 导出
# ============================================================================
def _export_pdf(file_id: str, target_lang: str, mode: str) -> BytesIO:
    """PDF 导出：在原 PDF 上叠加翻译文本。"""
    import fitz

    file_path = UPLOAD_DIR / file_id
    pages = _load_translated_pages(file_id, target_lang)

    doc = fitz.open(str(file_path))

    for page_data in pages:
        page_num = page_data.get("page", 0) - 1
        if page_num < 0 or page_num >= len(doc):
            continue

        page = doc[page_num]
        text = page_data.get("text", "")

        if mode == "translated":
            # 白底覆盖 + 新文本
            rect = page.rect
            page.draw_rect(rect, color=None, fill=(1, 1, 1))  # 白色覆盖
            page.insert_text(
                fitz.Point(50, 72),
                text,
                fontname="helv",
                fontsize=11,
            )
        elif mode == "bilingual":
            # 双栏：左原文右译文
            original_text = page.get_text("text") or ""
            half_w = page.rect.width / 2 - 20
            # 清空页面
            page.draw_rect(page.rect, color=None, fill=(1, 1, 1))
            page.insert_text(
                fitz.Point(30, 72),
                "[Original]\n" + original_text[:500],
                fontname="helv",
                fontsize=9,
            )
            page.insert_text(
                fitz.Point(half_w + 30, 72),
                f"[{target_lang.upper()}]\n" + text[:500],
                fontname="helv",
                fontsize=9,
            )

    buf = BytesIO()
    doc.save(buf)
    doc.close()
    buf.seek(0)
    logger.info("PDF 导出完成: %s → %s (%s)", file_id, target_lang, mode)
    return buf


# ============================================================================
# PPTX 导出
# ============================================================================
def _export_pptx(file_id: str, target_lang: str, mode: str) -> BytesIO:
    """PPTX 导出：将翻译文本按 shape 逐行回写。"""
    from pptx import Presentation

    file_path = UPLOAD_DIR / file_id
    pages = _load_translated_pages(file_id, target_lang)

    prs = Presentation(str(file_path))

    for page_data in pages:
        slide_num = page_data.get("page", 0) - 1
        if slide_num < 0 or slide_num >= len(prs.slides):
            continue

        slide = prs.slides[slide_num]
        text_lines = page_data.get("text", "").split("\n")
        line_idx = 0

        if mode == "translated":
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if line_idx < len(text_lines):
                            # 保留第一个 run 的样式，替换文本
                            if para.runs:
                                para.runs[0].text = text_lines[line_idx]
                                for extra in para.runs[1:]:
                                    extra.text = ""
                            else:
                                para.text = text_lines[line_idx]
                            line_idx += 1

        elif mode == "bilingual":
            # 交替：原文行 → 译文行
            original_lines: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            original_lines.append(t)

            for shape in slide.shapes:
                if shape.has_text_frame:
                    tf = shape.text_frame
                    tf.clear()
                    idx = 0
                    for orig in original_lines:
                        if idx < len(text_lines):
                            p = tf.add_paragraph()
                            p.text = f"[ORIG] {orig}"
                            p2 = tf.add_paragraph()
                            p2.text = f"[{target_lang.upper()}] {text_lines[idx]}"
                            idx += 1
                    break  # 只修改第一个 text_frame

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    logger.info("PPTX 导出完成: %s → %s (%s)", file_id, target_lang, mode)
    return buf


# ============================================================================
# Word 导出
# ============================================================================
def _export_word(file_id: str, target_lang: str, mode: str) -> BytesIO:
    """Word 导出：将翻译文本逐段回写，保留样式。"""
    from docx import Document

    file_path = UPLOAD_DIR / file_id
    pages = _load_translated_pages(file_id, target_lang)

    # 将分页的文本重新合并为行
    all_lines: list[str] = []
    for page_data in pages:
        text = page_data.get("text", "")
        all_lines.extend(text.split("\n"))

    doc = Document(str(file_path))
    line_idx = 0

    if mode == "translated":
        for para in doc.paragraphs:
            if line_idx < len(all_lines):
                if para.runs and all_lines[line_idx].strip():
                    para.runs[0].text = all_lines[line_idx]
                    for extra in para.runs[1:]:
                        extra.text = ""
                line_idx += 1

    elif mode == "bilingual":
        # 在每个原文段落后面插入译文
        original_paras = [p for p in doc.paragraphs if p.text.strip()]
        for i, orig in enumerate(original_paras):
            if i < len(all_lines):
                orig.text = f"[ORIG] {orig.text}"
                # 插入译文段落
                new_para = orig.insert_paragraph_after(f"[{target_lang.upper()}] {all_lines[i]}")

    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    logger.info("Word 导出完成: %s → %s (%s)", file_id, target_lang, mode)
    return buf
