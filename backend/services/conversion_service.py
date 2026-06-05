"""Document format conversion dispatcher.

Priority tiers:
  P0 — sync, <2s:   PDF→TXT, DOCX→TXT, PPTX→TXT
  P1 — sync ≤20MB:  PDF→DOCX, DOCX→PDF, PPTX→PDF, TXT→PDF, TXT→DOCX
  P2 — always async: PDF→PPTX

Windows: uses Office COM (Word/PowerPoint) for DOCX/PPTX → PDF (perfect fidelity).
"""

from __future__ import annotations

import io
import logging
import os
import tempfile
from pathlib import Path

from cache.file_cache import FileCache

logger = logging.getLogger("tri_doc.convert")

BASE_DIR = Path(__file__).resolve().parent.parent
UPLOAD_DIR = BASE_DIR / "uploads"

P0_FORMATS = {"txt"}
P2_PAIRS = {("pdf", "pptx")}

cache = FileCache()

SUPPORTED_CONVERSIONS: dict[tuple[str, str], str] = {
    ("pdf", "txt"): "P0",
    ("docx", "txt"): "P0",
    ("pptx", "txt"): "P0",
    ("pdf", "docx"): "P1",
    ("docx", "pdf"): "P1",
    ("pptx", "pdf"): "P1",
    ("txt", "pdf"): "P1",
    ("txt", "docx"): "P1",
    ("pdf", "pptx"): "P2",
}


def _source_ext(file_id: str) -> str:
    return Path(file_id).suffix.lower().lstrip(".")


def _read_source(file_id: str) -> bytes:
    path = UPLOAD_DIR / file_id
    if not path.exists():
        raise FileNotFoundError(f"Source file not found: {file_id}")
    return path.read_bytes()


def _file_size_mb(file_id: str) -> float:
    path = UPLOAD_DIR / file_id
    return path.stat().st_size / (1024 * 1024) if path.exists() else 0


# ---------------------------------------------------------------------------
# COM helper — converts any Office-supported format using the host app
# ---------------------------------------------------------------------------

def _com_convert(file_id: str, target_format: str, app_name: str, export_format, progress_callback=None):
    """Convert via Microsoft Office COM — perfect fidelity for DOCX/PPTX → PDF."""
    import pythoncom
    import win32com.client

    src = UPLOAD_DIR / file_id
    pythoncom.CoInitialize()

    app = None
    doc = None
    tmp_path = None

    try:
        if progress_callback:
            progress_callback(f"{app_name}启动", 0, 100)

        app = win32com.client.Dispatch(f"{app_name}.Application")
        app.Visible = False
        app.DisplayAlerts = 0  # wdAlertsNone / ppAlertsNone

        if progress_callback:
            progress_callback(f"{app_name}打开文档", 20, 100)

        doc = app.Documents.Open(str(src)) if app_name == "Word" else app.Presentations.Open(str(src), WithWindow=False)

        fd, tmp_path = tempfile.mkstemp(suffix=f".{target_format}")
        os.close(fd)

        if progress_callback:
            progress_callback(f"{app_name}导出", 40, 100)

        if app_name == "Word":
            doc.SaveAs2(tmp_path, FileFormat=export_format)
        else:
            doc.SaveAs(tmp_path, FileFormat=export_format)

        if progress_callback:
            progress_callback("读取结果", 90, 100)

        result = Path(tmp_path).read_bytes()
        return result

    finally:
        if doc is not None:
            try:
                doc.Close()
            except Exception:
                pass
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.unlink(tmp_path)
            except Exception:
                pass
        pythoncom.CoUninitialize()


# ---------------------------------------------------------------------------
# P0: Extract text (pure Python, fast)
# ---------------------------------------------------------------------------

def _pdf_to_txt(file_id: str) -> bytes:
    import fitz
    buf = io.BytesIO(_read_source(file_id))
    doc = fitz.open(stream=buf, filetype="pdf")
    parts = []
    for page in doc:
        parts.append(page.get_text())
    doc.close()
    return "\n\n".join(parts).encode("utf-8")


def _docx_to_txt(file_id: str) -> bytes:
    from docx import Document
    buf = io.BytesIO(_read_source(file_id))
    doc = Document(buf)
    parts = [p.text for p in doc.paragraphs]
    return "\n".join(parts).encode("utf-8")


def _pptx_to_txt(file_id: str) -> bytes:
    from pptx import Presentation
    buf = io.BytesIO(_read_source(file_id))
    prs = Presentation(buf)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        parts.append(t)
        parts.append("")
    return "\n".join(parts).encode("utf-8")


# ---------------------------------------------------------------------------
# P1: Office COM for DOCX/PPTX → PDF, pdf2docx for PDF → DOCX
# ---------------------------------------------------------------------------

def _pdf_to_docx(file_id: str, progress_callback=None) -> bytes:
    from pdf2docx import Converter
    src = UPLOAD_DIR / file_id
    buf = io.BytesIO()
    cv = Converter(str(src))
    cv.convert(buf)
    cv.close()
    return buf.getvalue()


def _docx_to_pdf(file_id: str, progress_callback=None) -> bytes:
    # wdFormatPDF = 17
    return _com_convert(file_id, "pdf", "Word", 17, progress_callback)


def _pptx_to_pdf(file_id: str, progress_callback=None) -> bytes:
    # ppSaveAsPDF = 32
    return _com_convert(file_id, "pdf", "PowerPoint", 32, progress_callback)


def _txt_to_pdf(file_id: str, progress_callback=None) -> bytes:
    import fitz
    text = _read_source(file_id).decode("utf-8", errors="replace")
    lines = text.split("\n")
    pdf = fitz.open()
    page = pdf.new_page(width=595, height=842)
    y = 72
    for line in lines:
        if y > 770:
            page = pdf.new_page(width=595, height=842)
            y = 72
        page.insert_text((72, y), line, fontsize=11)
        y += 14
    out = io.BytesIO()
    pdf.save(out)
    pdf.close()
    return out.getvalue()


def _txt_to_docx(file_id: str, progress_callback=None) -> bytes:
    from docx import Document
    text = _read_source(file_id).decode("utf-8", errors="replace")
    doc = Document()
    for line in text.split("\n"):
        doc.add_paragraph(line)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# P2: Image-based (no COM for PDF→PPTX since PowerPoint can't open PDF)
# ---------------------------------------------------------------------------

def _pdf_to_pptx(file_id: str, progress_callback=None) -> bytes:
    import fitz
    from pptx import Presentation
    from pptx.util import Inches

    buf = io.BytesIO(_read_source(file_id))
    pdf = fitz.open(stream=buf, filetype="pdf")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total = len(pdf)
    for i, page in enumerate(pdf):
        if progress_callback:
            progress_callback("PDF→PPTX (render)", i + 1, total)
        pix = page.get_pixmap(dpi=150)
        img_buf = io.BytesIO(pix.tobytes("png"))
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(img_buf, 0, 0, width=prs.slide_width, height=prs.slide_height)

    pdf.close()
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

_HANDLERS = {
    ("pdf", "txt"): _pdf_to_txt,
    ("docx", "txt"): _docx_to_txt,
    ("pptx", "txt"): _pptx_to_txt,
    ("pdf", "docx"): _pdf_to_docx,
    ("docx", "pdf"): _docx_to_pdf,
    ("pptx", "pdf"): _pptx_to_pdf,
    ("txt", "pdf"): _txt_to_pdf,
    ("txt", "docx"): _txt_to_docx,
    ("pdf", "pptx"): _pdf_to_pptx,
}

_EXT_MAP = {"pdf": ".pdf", "docx": ".docx", "pptx": ".pptx", "txt": ".txt"}


def convert_document(file_id: str, target_format: str, progress_callback=None) -> tuple[bytes, str]:
    src_ext = _source_ext(file_id)
    target_format = target_format.lower().strip()
    pair = (src_ext, target_format)

    if pair not in _HANDLERS:
        raise ValueError(f"Unsupported conversion: {src_ext} → {target_format}")

    cached = cache.get(file_id, target_format)
    if cached is not None:
        return cached.getvalue(), _EXT_MAP[target_format]

    handler = _HANDLERS[pair]
    data = handler(file_id, progress_callback=progress_callback) if pair in P2_PAIRS else handler(file_id)

    cache.put(file_id, target_format, data, _EXT_MAP[target_format])
    return data, _EXT_MAP[target_format]


def is_async(file_id: str, target_format: str) -> bool:
    src_ext = _source_ext(file_id)
    pair = (src_ext, target_format)
    if pair in P2_PAIRS:
        return True
    return _file_size_mb(file_id) > 20


def get_target_formats(source_ext: str) -> list[str]:
    targets = set()
    for (src, tgt) in SUPPORTED_CONVERSIONS:
        if src == source_ext.lower():
            targets.add(tgt)
    return sorted(targets)
