"""
TriDoc 文档解析服务 — 统一的结构化解析接口。

输出格式（统一 ContentBlock）:
  {
    "pages": [
      {
        "page": int,
        "text": str,             # 纯文本（用于编辑）
        "blocks": ContentBlock[] # 结构化（用于样式保留）
      }
    ]
  }

ContentBlock:
  {
    "type": "textbox" | "image" | "table",
    "position": {x, y, w, h} | null,
    "editable": bool,
    "paragraphs": [
      {
        "alignment": str | null,
        "runs": [{text, font_name, font_size, bold, italic, color}]
      }
    ]
  }
"""

from __future__ import annotations

import logging
import os
from pathlib import Path
from typing import Any

logger = logging.getLogger("tri_doc.parser")

BASE_DIR = Path(__file__).resolve().parent.parent
UPLOAD_DIR = BASE_DIR / "uploads"


# ============================================================================
# 统一解析入口
# ============================================================================
def parse_document(file_id: str) -> dict:
    """
    根据文件扩展名分派解析器，返回统一的页面结构。

    Returns:
        {"file_id": str, "file_type": "pdf"|"pptx"|"word", "pages": [...]}
    """
    file_path = UPLOAD_DIR / file_id
    if not file_path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = file_path.suffix.lower()
    if ext == ".pdf":
        return parse_pdf(file_id)
    elif ext == ".pptx":
        return parse_pptx(file_id)
    elif ext == ".docx":
        return parse_word(file_id)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


# ============================================================================
# PDF 解析 — PyMuPDF (fitz)
# ============================================================================
def parse_pdf(file_id: str) -> dict:
    """PDF 解析：提取每页文本 + 文本块位置信息。"""
    import fitz  # PyMuPDF

    file_path = UPLOAD_DIR / file_id
    doc = fitz.open(str(file_path))
    pages: list[dict] = []

    for i, page in enumerate(doc, 1):
        text = page.get_text("text") or ""
        blocks_raw = page.get_text("dict").get("blocks", [])

        content_blocks: list[dict] = []
        for block in blocks_raw:
            if block.get("type") == 0:  # 文本块
                cb = _pdf_text_block(block)
                if cb:
                    content_blocks.append(cb)
            elif block.get("type") == 1:  # 图片块
                bbox = block.get("bbox", [0, 0, 0, 0])
                content_blocks.append({
                    "type": "image",
                    "position": {"x": bbox[0], "y": bbox[1], "w": bbox[2] - bbox[0], "h": bbox[3] - bbox[1]},
                    "editable": False,
                    "paragraphs": [],
                })

        pages.append({
            "page": i,
            "text": text.strip(),
            "blocks": content_blocks,
        })

    doc.close()
    logger.info("PDF 解析完成: %s → %d 页", file_id, len(pages))
    return {"file_id": file_id, "file_type": "pdf", "pages": pages}


def _pdf_text_block(block: dict) -> dict | None:
    """将 PyMuPDF text block 转为 ContentBlock。"""
    lines = block.get("lines", [])
    if not lines:
        return None

    all_spans: list[dict] = []
    for line in lines:
        for span in line.get("spans", []):
            all_spans.append(span)

    if not all_spans:
        return None

    # 取第一个 span 的样式作为段落样式
    first = all_spans[0]
    font_name = first.get("font", "Helvetica")
    font_size = round(first.get("size", 10), 1)
    bbox = block.get("bbox", [0, 0, 0, 0])

    runs: list[dict] = []
    for span in all_spans:
        color_int = span.get("color", 0)
        color_hex = f"#{color_int:06x}" if isinstance(color_int, int) else "#000000"
        flags = span.get("flags", 0)
        runs.append({
            "text": span.get("text", ""),
            "font_name": span.get("font", font_name),
            "font_size": round(span.get("size", font_size), 1),
            "bold": bool(flags & 2**3),
            "italic": bool(flags & 2**1),
            "color": color_hex,
        })

    return {
        "type": "textbox",
        "position": {"x": bbox[0], "y": bbox[1], "w": bbox[2] - bbox[0], "h": bbox[3] - bbox[1]},
        "editable": True,
        "paragraphs": [{"alignment": None, "runs": runs}],
    }


# ============================================================================
# PPTX 解析 — python-pptx
# ============================================================================
def parse_pptx(file_id: str) -> dict:
    """PPTX 解析：每页 slide → shapes → text_frame → paragraphs → runs。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    file_path = UPLOAD_DIR / file_id
    prs = Presentation(str(file_path))
    slides: list[dict] = []

    for i, slide in enumerate(prs.slides, 1):
        all_text: list[str] = []
        content_blocks: list[dict] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                cb = _pptx_text_block(shape)
                if cb:
                    content_blocks.append(cb)
                    # 收集纯文本
                    for p in cb["paragraphs"]:
                        txt = "".join(r["text"] for r in p.get("runs", []))
                        if txt.strip():
                            all_text.append(txt.strip())
            elif shape.shape_type == 13:  # Picture
                content_blocks.append({
                    "type": "image",
                    "position": {
                        "x": shape.left, "y": shape.top,
                        "w": shape.width, "h": shape.height,
                    },
                    "editable": False,
                    "paragraphs": [],
                })

        slides.append({
            "page": i,
            "text": "\n".join(all_text),
            "blocks": content_blocks,
        })

    logger.info("PPTX 解析完成: %s → %d 页", file_id, len(slides))
    return {"file_id": file_id, "file_type": "pptx", "pages": slides}


def _pptx_text_block(shape) -> dict | None:
    """将 PPTX shape 转为 ContentBlock。"""
    tf = shape.text_frame
    paragraphs: list[dict] = []

    for para in tf.paragraphs:
        runs: list[dict] = []
        for run in para.runs:
            font = run.font
            runs.append({
                "text": run.text,
                "font_name": font.name or "Arial",
                "font_size": font.size / 12700 if font.size else 12,  # EMU → pt
                "bold": font.bold or False,
                "italic": font.italic or False,
                "color": str(font.color.rgb) if font.color and font.color.rgb else "#000000",
            })
        if runs:
            paragraphs.append({
                "alignment": str(para.alignment) if para.alignment else None,
                "runs": runs,
            })

    if not paragraphs:
        return None

    return {
        "type": "textbox",
        "position": {"x": shape.left, "y": shape.top, "w": shape.width, "h": shape.height},
        "editable": True,
        "paragraphs": paragraphs,
    }


# ============================================================================
# Word 解析 — python-docx
# ============================================================================
def parse_word(file_id: str) -> dict:
    """Word 解析：paragraphs → runs（保留样式），按段落分页（每10段≈1页）。"""
    from docx import Document

    file_path = UPLOAD_DIR / file_id
    doc = Document(str(file_path))

    all_paragraphs: list[dict] = []
    all_text: list[str] = []

    for para in doc.paragraphs:
        runs: list[dict] = []
        for run in para.runs:
            font = run.font
            runs.append({
                "text": run.text,
                "font_name": font.name or "Calibri",
                "font_size": font.size / 12700 if font.size else 11,
                "bold": font.bold or False,
                "italic": font.italic or False,
                "color": str(font.color.rgb) if font.color and font.color.rgb else "#000000",
            })
        para_dict = {
            "alignment": str(para.alignment) if para.alignment else None,
            "runs": runs,
        }
        all_paragraphs.append(para_dict)
        txt = "".join(r["text"] for r in runs)
        if txt.strip():
            all_text.append(txt.strip())

    # 按每 15 个段落分组为"页"（Word 无天然分页）
    PARS_PER_PAGE = 15
    pages: list[dict] = []
    for i in range(0, len(all_paragraphs), PARS_PER_PAGE):
        chunk = all_paragraphs[i : i + PARS_PER_PAGE]
        page_num = i // PARS_PER_PAGE + 1
        pages.append({
            "page": page_num,
            "text": "\n".join(
                "".join(r["text"] for r in p["runs"]) for p in chunk
            ),
            "blocks": [{
                "type": "textbox",
                "position": None,
                "editable": True,
                "paragraphs": chunk,
            }],
        })

    logger.info("Word 解析完成: %s → %d 段 / %d 页", file_id, len(all_paragraphs), len(pages))
    return {"file_id": file_id, "file_type": "word", "pages": pages}
